import streamlit as st
import base64
import pptx
from pptx.util import Inches, Pt
import os
from ctransformers import AutoModelForCausalLM
from dotenv import load_dotenv
import pandas as pd
import numpy as np
from textwrap import dedent
from typing import List
from fastembed import TextEmbedding
from groq import Groq
from groq.types.chat import ChatCompletionMessage
import streamlit as st
import re
load_dotenv()


# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)



# local_llm = r'llama-2-7b-chat.Q4_K_M.gguf'

# llm = AutoModelForCausalLM.from_pretrained(local_llm, model_type='llama')


# Set up the Groq API key and client
os.environ["GROQ_API_KEY"] = "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"
client = Groq()

MODEL = "llama-3.1-70b-versatile"
# MODEL = "7b-chat-q4_K_M"

TEMPERATURE = 0
MAX_TOKENS = 4096

def call_model(query: str, messages: List) -> ChatCompletionMessage:
    messages.append(
        {
            "role": "user",
            "content": query,
        },
    )
    response = client.chat.completions.create(
        model=MODEL,
        messages=messages,
        temperature=TEMPERATURE,
        max_tokens=MAX_TOKENS,
    )
    message = response.choices[0].message
    messages.append(message)
    return message


# SYSTEM_PROMPT = '''You are an AI-Powered Presentation Generator, designed to generate detailed professional PowerPoint presentations with complete explanation of the topics and the topic that user provide must divide in slide by slide sections.
# Your primary task is to generate content based on user inputs related to presentation slides. You're always helpful to generate  presenetation only based on the provided topic. If you don't unserstand the topic - just reply with an excuse that you
# don't know. Keep your content brief and detailed. Be kind and respectful.'''
SYSTEM_PROMPT = ''


def generate_slide_titles(topic, no_of_slides):
    prompt = f"Generate {no_of_slides} slide titles for the topic '{topic}'."

    # response = llm(prompt)


    if prompt:
        messages = [
            {"role": "system", "content": SYSTEM_PROMPT},
        ]

        # with st.spinner("Generating response..."):
        bot_response = call_model(prompt, messages)

        # st.write("**Bot Response:**")
        # st.write(bot_response.content)

    
    return bot_response.content.split("\n")

def generate_slide_content(slide_title, presentation_category):
    prompt = f"Generate content for the slide: '{slide_title}' for the category {presentation_category}. dont give 'Image Suggestions', 'Color Scheme', 'Fonts', and 'Slide Title'"
    # print(prompt)
    # response = llm(prompt)
    # print(f'\nresponse of the title {slide_title}\n', response)
    if prompt:
        messages = [
            {"role": "system", "content": SYSTEM_PROMPT},
        ]

        # with st.spinner("Generating response..."):
        bot_response = call_model(prompt, messages)
    print('\n\nbot_response:\n')
    print(bot_response.content)

    return bot_response.content.replace('**', '').replace('*', '').replace('• ', '')


def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        # Customize font size for titles and content
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    # prs.save(f"generated_ppt/{topic}_presentation.pptx")
    prs.save(f"generated_ppt/3_presentation.pptx")


####################################################################################################################################################################################    

#                                                                          Streamlit UI integration to Application                                                       

######################################################################################################################################################################################
def main():
    st.title("Presentation Generator")
    # st.subheader("Empowering Presentations with Open-source AI equipped PowerPoint Generator")
    st.markdown('<style>h1{color: orange; text-align: center;}</style>', unsafe_allow_html=True)
    st.markdown('<style>h3{color: pink; text-align: center;}</style>', unsafe_allow_html=True)

    topic = st.text_input("Enter the topic for your presentation:")
    # topic = re.sub(r'[^a-zA-Z]', '', topic)
    topic = re.sub(r'[^a-zA-Z\s]', '', topic)
    # presentation_category = st.text_input("Enter presentation category:")
    presentation_category = st.selectbox('Select category',('Business presentations', 'Motivational presentation', 'Persuasive presentations', 'Demonstrative presentations', 'Informative', 'Storytelling', 'Academic presentations', 'Sales presentation', 'Instructional', 'Progress presentations', 'Coach style', 'Connector style', 'Decision-making presentations', 'Instructor style', 'Problem solving'))
    no_of_slides = st.number_input("No. of slides:", min_value=3, max_value=10, step=1)
    generate_button = st.button("Generate Presentation")

    if generate_button and topic and no_of_slides and presentation_category:
        st.info("Generating presentation... Please wait.")
        # heading = st.text_input("Give a heading:")
        # description = st.text_input("Give an idea the of the description:")

        slide_titles = generate_slide_titles(topic, no_of_slides)
        filtered_slide_titles = [item.replace('**', '') for item in slide_titles if item.strip() != '']
        filtered_slide_titles = [title for title in filtered_slide_titles if title[0].isdigit()]
        filtered_slide_titles = [title.split('. ', 1)[1] for title in filtered_slide_titles if title[0].isdigit()]
        
        # filtered_slide_titles = filtered_slide_titles[1:]
        print(f"\n\nSlide Title: {filtered_slide_titles}\n\n" )

        slide_contents = [generate_slide_content(title, presentation_category) for title in filtered_slide_titles[1:]]
        print("\n\n\n Slide Contents: ", slide_contents, '\n\n')
        create_presentation(topic, filtered_slide_titles, slide_contents)
        print("Presentation generated successfully!")

        st.success("Presentation generated successfully!")
        st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)

def get_ppt_download_link(topic):
    # Remove invalid characters from the filename
    # pattern = r'[^A-Za-z0-9 ]'
    # sanitized_topic = re.sub(pattern, '', topic)
    ppt_filename = f"generated_ppt/3_presentation.pptx"
    print('ppt_filenameppt_filenameppt_filenameppt_filename: ', ppt_filename)
    
    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'

if __name__ == "__main__":
    main()