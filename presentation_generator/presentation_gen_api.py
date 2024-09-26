from fastapi import FastAPI, File, UploadFile, Form, HTTPException, UploadFile, File, Body
from fastapi.responses import JSONResponse
import streamlit as st
import base64
import pptx
from pptx.util import Inches, Pt
import os
from ctransformers import AutoModelForCausalLM
from dotenv import load_dotenv
import pandas as pd
import numpy as np
from textwrap import dedent
from typing import List
from fastembed import TextEmbedding
from groq import Groq
from groq.types.chat import ChatCompletionMessage
import streamlit as st
import re
load_dotenv()


# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)



# local_llm = r'llama-2-7b-chat.Q4_K_M.gguf'

# llm = AutoModelForCausalLM.from_pretrained(local_llm, model_type='llama')


# Set up the Groq API key and client
os.environ["GROQ_API_KEY"] = "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"
client = Groq()

MODEL = "llama-3.1-70b-versatile"
# MODEL = "7b-chat-q4_K_M"

TEMPERATURE = 0
MAX_TOKENS = 4096
app = FastAPI()

def call_model(query: str, messages: List) -> ChatCompletionMessage:
    messages.append(
        {
            "role": "user",
            "content": query,
        },
    )
    response = client.chat.completions.create(
        model=MODEL,
        messages=messages,
        temperature=TEMPERATURE,
        max_tokens=MAX_TOKENS,
    )
    message = response.choices[0].message
    messages.append(message)
    return message


SYSTEM_PROMPT = '''You are an AI-Powered Presentation Generator, designed to generate 2 best detailed version of professional PowerPoint presentations with complete explanation of the topics and the topic that user provide must divide in slide by slide sections.
Your primary task is to generate content based on user inputs related to presentation slides. You're always helpful to generate  presenetation only based on the provided topic. If you don't unserstand the topic - just reply with an excuse that you
don't know. Keep your content brief and detailed. Be kind and respectful.'''

def generate_slide_titles(topic, no_of_slides, presentation_category, slide_number):
    prompt = f"Generate three different types of title for slide number '{slide_number}' for the topic '{topic}' with the category {presentation_category}."
    if prompt:
        messages = [
            {"role": "system", "content": SYSTEM_PROMPT},
        ]

        bot_response = call_model(prompt, messages)


    
    return bot_response.content.split("\n")

def generate_slide_content(topic, title, presentation_category, step2_description):
    prompt = f"Generate 2 best versions of the content for the slide: '{title}' with the category {presentation_category}, description {step2_description}, and topic {topic}. remember to just only give 3 lines related to the detail and 3-4 bullet points and dont give 'Image Suggestions', 'Color Scheme', 'Fonts', 'Key Points to Discuss',  and 'Slide Title just give detail and keypoints"

    
    if prompt:
        messages = [
            {"role": "system", "content": SYSTEM_PROMPT},
        ]

        bot_response = call_model(prompt, messages)

    return bot_response.content.replace('**', '').replace('*', '').replace('• ', '').replace('Detail:', '').replace('Detail: \n', '').replace('Key Points:\n', '')

def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        # Customize font size for titles and content
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    # prs.save(f"generated_ppt/{topic}_presentation.pptx")
    prs.save(f"api_test_generated_presentation.pptx")


@app.post("/title_gen")
async def title(topic: str = Form(...), no_of_slides: str = Form(...), selected_category: str = Form(...), slide_number: int = Form(...)):
    
    try:
        if no_of_slides == slide_number:
            titles = generate_slide_titles(topic, no_of_slides, selected_category, 'second last')
        else:
            titles = generate_slide_titles(topic, no_of_slides, selected_category, slide_number)
        # titles = generate_slide_titles(topic, no_of_slides, selected_category, slide_number)
        filtered_slide_titles = [item.replace('**', '') for item in titles if item.strip() != '']
        filtered_slide_titles = [title for title in filtered_slide_titles if title[0].isdigit()]
        filtered_slide_titles = [title.split('. ', 1)[1] for title in filtered_slide_titles if title[0].isdigit()]
        filtered_slide_titles = [title.split(': ')[1] if len(title.split(': ')) > 1 else title for title in filtered_slide_titles ]
        filtered_slide_titles = [item.replace('"', '') for item in filtered_slide_titles]
        response_data = {
                "success": True,
                "responseMessage": filtered_slide_titles,
                "responseCode": "200",
            }
    except Exception as e:
        response_data = {
            "success": False,
            "responseMessage": f"Error: {str(e)}",
            "responseCode": "500",
            "data": {}
        }

    return JSONResponse(content=response_data)


@app.post("/content_gen")
async def content(topic: str = Form(...), selected_slide_title: str = Form(...), selected_category: str = Form(...), des: str = Form(...)):
    
    try:
        text_versions = generate_slide_content(topic, selected_slide_title, selected_category, des)
        text_versions = text_versions.replace('Version 1:', '').split('Version 2:')
        response_data = {
                "success": True,
                "responseMessage": text_versions,
                "responseCode": "200",
            }
    except Exception as e:
        response_data = {
            "success": False,
            "responseMessage": f"Error: {str(e)}",
            "responseCode": "500",
            "data": {}
        }

    return JSONResponse(content=response_data)

@app.post("/create_ppt")
async def content(topic: str = Form(...), presen_titles: List[str] = Body(...), presen_content: List[str] = Body(...)):
    presen_titles = [item.strip() for item in presen_titles[0].split(',')]
    # print('presen_titles: ',presen_titles, 'type presen_titles: ', type(presen_titles))
    presen_content = [item.strip() for item in presen_content[0].replace('\\n','\n').split(',\n\n')]
    # print('\n\npresen_content: ', presen_content)
    # print('presen_content: ',presen_content, 'type presen_content: ', type(presen_content), 'length: ', len(presen_content))

    # print('presen_titles: ', new_list, 'type presen_titles: ', type(new_list), 'length: ', len(new_list))

    try:
        create_presentation(topic, presen_titles, presen_content)
        response_data = {
                "success": True,
                "responseMessage": 'presentation_created',
                "responseCode": "200",
            }
    except Exception as e:
        response_data = {
            "success": False,
            "responseMessage": f"Error: {str(e)}",
            "responseCode": "500",
            "data": {}
        }

    return JSONResponse(content=response_data)    



# ['"\\n\\n \\nArtificial Intelligence (AI) is a branch of computer science that focuses on creating intelligent machines capable of performing tasks that typically require human intelligence. AI combines computer science, mathematics, and engineering to develop algorithms and statistical models that enable machines to learn, reason, and interact with their environment. The field of AI has been rapidly evolving, with significant advancements in recent years.\\n\\nAI involves the development of algorithms and statistical models that enable machines to learn from data.\\nAI systems can be classified into two main categories: narrow or weak AI, and general or strong AI.\\nNarrow AI is designed to perform a specific task, whereas general AI aims to match human intelligence and capabilities.\\nAI has numerous applications across various industries, including healthcare, finance, and transportation.\\n\\n', '\\nArtificial Intelligence (AI) is a multidisciplinary field that seeks to create machines capable of simulating human intelligence, learning, and problem-solving abilities. AI draws inspiration from human cognition, neuroscience, and computer science to develop intelligent systems that can interact with their environment and adapt to new situations. The study of AI has led to significant breakthroughs in areas such as machine learning, natural language processing, and computer vision.\\n\\nAI is a multidisciplinary field that combines insights from computer science, neuroscience, and cognitive psychology.\\nAI systems can be trained using various machine learning algorithms, including supervised, unsupervised, and reinforcement learning.\\nAI has numerous applications in areas such as image recognition, speech recognition, and natural language processing.\\nAI has the potential to transform industries and revolutionize the way we live and work.']
# '\n\n Artificial Intelligence (AI) has revolutionized the business world by transforming the way companies operate, make decisions, and interact with customers. AI refers to the development of computer systems that can perform tasks that typically require human intelligence, such as learning, problem-solving, and decision-making. As AI continues to evolve, it is essential for businesses to understand its potential and limitations.\n\nAI is a subset of computer science that focuses on creating intelligent machines.\nAI has numerous applications in business, including customer service, marketing, and finance.\nAI can help businesses automate repetitive tasks, improve efficiency, and reduce costs.\nAI can also enhance decision-making by providing insights and predictions based on data analysis.\n\n', '\n\n Artificial Intelligence (AI) has revolutionized the business world, transforming the way companies operate and make decisions. Here are the key takeaways from our discussion on AI.\n\nAI can automate repetitive tasks, freeing up human resources for more strategic and creative work.\nAI-powered analytics can provide valuable insights, enabling businesses to make data-driven decisions.\nAI can enhance customer experience through personalized interactions and real-time support.\nAI adoption can lead to increased efficiency, productivity, and competitiveness.\n\n'